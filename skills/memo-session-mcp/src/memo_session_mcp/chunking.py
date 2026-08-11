"""Markdown chunking utilities."""

from __future__ import annotations

import hashlib
import re
from dataclasses import dataclass
from pathlib import Path


@dataclass
class Chunk:
    content: str
    section: str
    chunk_index: int
    line_start: int


_HEADING = re.compile(r"^(#{1,4})\s+(.+)$", re.MULTILINE)


def file_sha256(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as fh:
        for block in iter(lambda: fh.read(65536), b""):
            h.update(block)
    return h.hexdigest()


def chunk_markdown(
    text: str,
    *,
    max_chars: int = 1800,
    overlap: int = 200,
) -> list[Chunk]:
    sections: list[tuple[str, str, int]] = []
    lines = text.splitlines()
    current_title = "(root)"
    start_line = 1
    buf: list[str] = []

    def flush(title: str, start: int) -> None:
        body = "\n".join(buf).strip()
        if body:
            sections.append((title, body, start))

    for i, line in enumerate(lines, start=1):
        m = re.match(r"^(#{1,4})\s+(.+)$", line)
        if m and buf:
            flush(current_title, start_line)
            buf = [line]
            current_title = m.group(2).strip()
            start_line = i
        else:
            if not buf:
                start_line = i
            buf.append(line)
    flush(current_title, start_line)

    chunks: list[Chunk] = []
    idx = 0
    for title, body, line_start in sections:
        if len(body) <= max_chars:
            chunks.append(Chunk(body, title, idx, line_start))
            idx += 1
            continue
        pos = 0
        while pos < len(body):
            piece = body[pos : pos + max_chars]
            chunks.append(Chunk(piece.strip(), title, idx, line_start))
            idx += 1
            if pos + max_chars >= len(body):
                break
            pos += max_chars - overlap
    return chunks


def extract_pdf_text(path: Path) -> str:
    try:
        import fitz  # pymupdf
    except ImportError as exc:
        raise RuntimeError(
            "PDF support requires: pip install memo-session-mcp[documents]"
        ) from exc
    doc = fitz.open(path)
    parts: list[str] = []
    for page in doc:
        parts.append(page.get_text())
    return "\n".join(parts)


def extract_xlsx_text(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError(
            "Excel support requires: pip install memo-session-mcp[documents]"
        ) from exc
    wb = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"## Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def read_document(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".md", ".txt", ".markdown"}:
        return path.read_text(encoding="utf-8", errors="replace")
    if suffix == ".pdf":
        return extract_pdf_text(path)
    if suffix in {".xlsx", ".xlsm"}:
        return extract_xlsx_text(path)
    if suffix == ".pptx":
        return _pptx_text(path)
    return ""


def _pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError as exc:
        raise RuntimeError(
            "PPTX support requires: pip install python-pptx (memo-session-mcp[documents])"
        ) from exc
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)
